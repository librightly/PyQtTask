# coding: utf-8

import os
#import datetime
from LDatetime import LDateTime
from File import File
from SEProblem import SEProblem
from SEProblemSet import SEProblemSet

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from PyQt5.QtWidgets import QMessageBox

class SESheetPPT(File):
    """问题单PPT类"""
    def __init__(self):
        super().__init__()
        
        """问题单内包含的问题"""
        self.problem_set = SEProblemSet()
        """问题单项目ID"""
        self.project_id = ""
        """问题单地址"""
        #self.filename = ""
        """问题单格式版本号"""
        self.version = "V0.0"
        
        """是否已读写标记"""
        self.isread = False 
        self.iswrite = False 
    
    """读取问题单PPT"""
    def read_problem_sheet(self,proressBar,strErrorLog):
        prs = Presentation(self.filename)
        slide_layouts = prs.slide_layouts
        #设置进度条
        proressBar.setMaximum(len(prs.slides))
        proressBar.setValue(1)
        
        #读取问题单首页信息
        head_slide = prs.slides[0]
        #QMessageBox.information(None,'消息',str(len(prs.slides)),QMessageBox.Yes | QMessageBox.No)
        '''if len(prs.slides) ！= 3:
            strErrorLog += "SE问题PPT《"
            strErrorLog += strPPTName
            strErrorLog += "》首页格式有误！\n"'''
        for shape in head_slide.shapes:
            if shape.has_text_frame:
                strPPTName = str(shape.text)                
                if strPPTName.find("SE问题单") != -1:
                    strPPTNameList = strPPTName.split("-",)
                    if len(strPPTNameList) == 3:
                        self.project_id = strPPTNameList[0]
                        self.name = strPPTNameList[1]
                    else:
                        strErrorLog = "SE问题PPT《" + strPPTName + "》文件名格式有误！\n"
                        QMessageBox.information(None,'消息',strErrorLog,QMessageBox.Yes | QMessageBox.No)
                        return False
                elif len(strPPTName) == 4:
                    self.version = strPPTName
                    continue
                else:
                    continue
            else:
                continue

        #self.project_id = self.project_id[0,self.project_id.count()-2]        
        #遍历Slides
        for i in range(1,len(prs.slides)):
            slide = prs.slides[i]
            #获取Slide版式id
            #slide_layout = slide.slide_layout
            index = slide_layouts.index(slide.slide_layout)
            #判断Slide是否是问题单页版式
            if index >= 1 and index <= 3:
                #读取问题单页信息
                problem = SEProblem()
                #将此问题放置在问题集合尾端
                self.problem_set.append(problem)
                #遍历Slide内Shapes
                for j in range(0,len(slide.shapes)):
                    shape = slide.shapes[j]
                    if shape.has_table:  #读取表格内文字内容
                        if self.version == "V1.0":
                            table = shape.table
                            problem.project = self.project_id
                            problem.index = self.name + '-' + table.cell(1,0).text
                            problem.creater = table.cell(1,1).text
                            problem.area_product = table.cell(1,3).text
                            problem.area_division = table.cell(1,5).text
                            #problem.source = table.cell(1,7).text
                            #problem.index_source = table.cell(1,9).text
                            #problem.final_solution = table.cell(1,10).text
                            problem.brief_and_reason = table.cell(3,0).text
                            problem.suggestion = table.cell(6,0).text
                            problem.detail_and_progress = table.cell(3,4).text
                            problem.prevention_and_result = table.cell(11,4).text
                            problem.status = table.cell(7,11).text
                            problem.stakeholder = table.cell(9,10).text
                            #problem.top = table.cell(3,10).text
                            problem.version_start_1 = table.cell(13,4).text
                            problem.version_end_1 = table.cell(13,8).text
                            problem.version_start_2 = table.cell(14,4).text
                            problem.version_end_2 = table.cell(14,8).text
                            problem.version_start_3 = table.cell(15,4).text
                            problem.version_end_3 = table.cell(15,8).text
                            problem.version_start_4 = table.cell(16,4).text
                            problem.version_end_4 = table.cell(16,8).text
                            #problem.version_start_5 = table.cell(17,4).text
                            #problem.version_end_5 = table.cell(17,8).text
                            problem.note = table.cell(18,6).text
                            problem.risk = index

                            #获取表格内时间
                            problem.date_created.fromString(table.cell(1,2).text,True)
                            problem.date_planed.fromString(table.cell(11,10).text,True)
                            #problem.date_followed.fromString(table.cell(13,10).text,False)
                            problem.date_closed.fromString(table.cell(16,10).text,False)
                        elif self.version == "V2.0":
                            table = shape.table
                            problem.project = self.project_id
                            problem.index = self.name + '-' + table.cell(1,0).text
                            problem.creater = table.cell(1,1).text
                            problem.area_product = table.cell(1,3).text
                            problem.area_division = table.cell(1,5).text
                            problem.MIR_id = table.cell(1,7).text
                            problem.MR_id = table.cell(1,11).text
                            problem.style_SE = table.cell(1,13).text
                            problem.brief_and_reason = table.cell(3,0).text
                            problem.suggestion = table.cell(6,0).text
                            problem.detail_and_progress = table.cell(3,4).text
                            problem.prevention_and_result = table.cell(10,4).text
                            problem.stakeholder = table.cell(8,13).text
                            problem.status = table.cell(3,13).text
                            problem.version_start_1 = table.cell(12,4).text
                            problem.version_end_1 = table.cell(12,8).text
                            problem.version_start_2 = table.cell(13,4).text
                            problem.version_end_2 = table.cell(13,8).text
                            problem.version_start_3 = table.cell(14,4).text
                            problem.version_end_3 = table.cell(14,8).text
                            problem.version_start_4 = table.cell(15,4).text
                            problem.version_end_4 = table.cell(15,8).text
                            problem.LLR = table.cell(16,6).text
                            problem.stlye_LLR = table.cell(16,10).text
                            problem.update_MR = table.cell(16,14).text
                            problem.EPorPPV = table.cell(15,13).text
                            problem.note = table.cell(17,6).text

                            problem.risk = index

                            #获取表格内时间
                            problem.date_created.fromString(table.cell(1,2).text,True)
                            problem.date_planed.fromString(table.cell(10,13).text,True)
                            problem.date_closed.fromString(table.cell(12,13).text,False)
                        else:
                            strErrorLog = "SE问题PPT《" + strPPTName + "》问题单版本号有误！\n"
                            QMessageBox.information(None,'消息',strErrorLog,QMessageBox.Yes | QMessageBox.No)
                            return False
            else: continue
            proressBar.setValue(i)
        proressBar.setValue(len(prs.slides))
        self.isread = True
        
        #计算问题总数
        self.problem_set.analysis(strErrorLog)
        return True
